# ＨＴＭＬからの情報収集
# エラーが出たら、pip install lxml
import pandas as pd
KTQurl = 'https://www.city.kitakyushu.lg.jp/ho-huku/18901209.html'
KTQ_per100K7 = pd.read_html(KTQurl, header = 0)[1].iat[1,1]

# ＣＳＶからの情報収集
NHKurl = 'https://raw.githubusercontent.com/u-10bei/covid-19_JPdata/main/data/COVID-19_NHK.csv'
NHK = pd.read_csv(NHKurl).iloc[:,[0,4]]

# データ整形
NHK.reset_index(inplace=True)
NHK['Id'] = NHK['index']+1
NHK['都道府県'] = NHK['Id'].apply(lambda _: str(_))+'.'+NHK['Pref']
NHK['人'] = NHK['per100K7'].round(2)
NHK_chart = NHK.iloc[0:20,[4,5]]

# パワーポイントの操作用
# エラーが出たら、pip install python-pptx
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt

# 和暦の日付書式設定用
from datetime import date, timedelta
yesterday = date.today() - timedelta(days=1)
reiwa = 'R'+str(yesterday.year%100-18).zfill(2)
reiwa2 = '令和'+str(yesterday.year%100-18)+'年'
wareki = reiwa + yesterday.strftime('%m%d')
wareki2 = reiwa2 + str(yesterday.month) + '月' + str(yesterday.day) + '日現在'

# pptxの呼び出し
ppturl = './hinagata.pptx'
prs = Presentation(ppturl)
sld = prs.slides[0]

# テキストボックスの編集
text1 = sld.shapes[3].text_frame.paragraphs[0]
text1.text = wareki2
text1.font.bold = True
text1.font.size = Pt(14)
text2 = sld.shapes[4].text_frame.paragraphs[0]
text2.text = '※北九州市　'+KTQ_per100K7
text2.font.size = Pt(16)

# グラフの編集
chart = sld.shapes[0].chart
chart_data = CategoryChartData()
chart_data.categories = NHK_chart['都道府県']
chart_data.add_series('人',(NHK_chart['人']),'#,##0.00')
chart.replace_data(chart_data)

# スライドの保存
outurl = '06_' + wareki + '_感染者数（10万人あたり）.pptx'
prs.save(outurl)

# gmail送信準備
msg_title = '１０万人あたり感染者数の資料について'
msg_text = 'おはようございます。\n表題の件について、添付ファイルのとおり送付いたします。'

# メールサーバ接続　アカウントとかパスとかは、自分で作ってね
import smtplib
from util import credential as cred
from email.mime.text import MIMEText
from email.mime.multipart import MIMEMultipart
from email.mime.application import MIMEApplication
from email.utils import formatdate
#SMTPのオブジェクト作成。GmailのSMTPポートは587
smtpobj = smtplib.SMTP('smtp.gmail.com', 587)

#メールサーバに対する応答
smtpobj.ehlo()
#暗号化通信開始
smtpobj.starttls()
smtpobj.ehlo()
#ログイン
smtpobj.login(cred.GOOGLE_ACCOUNT, cred.GOOGLE_PASSWORD)

# メッセージとかの編集
send_list = cred.MAIL_TO.split(',')
msg = MIMEMultipart()
msg['Subject'] = msg_title
msg['From'] = cred.GOOGLE_ACCOUNT
msg['To'] = cred.MAIL_TO
msg['Date'] = formatdate(localtime=True)
msg.attach(MIMEText(msg_text))

with open(outurl,'rb') as f:
  attachment = MIMEApplication(f.read())

attachment.add_header('Content-Disposition','attachment',filename=outurl)
msg.attach(attachment)

smtpobj.sendmail(cred.GOOGLE_ACCOUNT,send_list,msg.as_string())
smtpobj.quit()